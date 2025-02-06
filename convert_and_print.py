import os
import tempfile
import win32api
import win32print
from tkinter import Tk, filedialog
from pptx import Presentation
from fpdf import FPDF

def select_pptx_file():
    Tk().withdraw()
    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
    return file_path

def convert_pptx_to_pdf(pptx_path, pdf_path):
    prs = Presentation(pptx_path)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    for i in range(0, len(prs.slides), 2):
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt=f"Slide {i+1}", ln=True, align='C')
        
        if i+1 < len(prs.slides):
            pdf.cell(200, 10, txt=f"Slide {i+2}", ln=True, align='C')
    
    pdf.output(pdf_path)

def print_pdf(pdf_path):
    printer_name = win32print.GetDefaultPrinter()
    win32api.ShellExecute(0, "print", pdf_path, f'"{printer_name}"', ".", 0)

def main():
    pptx_file = select_pptx_file()
    if not pptx_file:
        print("No file selected. Exiting.")
        return
    
    temp_dir = tempfile.gettempdir()
    pdf_file = os.path.join(temp_dir, "converted_slides.pdf")
    
    print("Converting PPTX to PDF...")
    convert_pptx_to_pdf(pptx_file, pdf_file)
    print("Conversion complete. Sending to printer...")
    print_pdf(pdf_file)
    print("Print job sent successfully!")

if __name__ == "__main__":
    main()

